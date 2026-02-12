from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define colors
GREEN = RGBColor(16, 185, 129)
DARK_BLUE = RGBColor(30, 41, 59)
LIGHT_GRAY = RGBColor(241, 245, 249)
WHITE = RGBColor(255, 255, 255)
TEXT_COLOR = RGBColor(15, 23, 42)

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    p = subtitle_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = GREEN
    
    return slide

def add_content_slide(prs, title, content_points):
    """Add a content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # Add colored header bar
    header_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1))
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = DARK_BLUE
    header_shape.line.color.rgb = GREEN
    header_shape.line.width = Pt(3)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Add content
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, point in enumerate(content_points):
        if i > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[i]
        p.text = point
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_COLOR
        p.space_before = Pt(8)
        p.space_after = Pt(8)
        p.level = 0
    
    return slide

# Slide 1: Title
add_title_slide(prs, "FACTOSIEM", "Data Ingestion, Visualization & Features Architecture")

# Slide 2: System Overview
add_content_slide(prs, "System Overview", [
    "🎯 Enterprise Threat Intelligence Platform",
    "🔄 Real-time data ingestion from multiple sources",
    "📊 Advanced visualization and analytics",
    "⚡ AI-Native investigation capabilities",
    "🛡️ Comprehensive threat detection and response"
])

# Slide 3: Data Ingestion
add_content_slide(prs, "Data Ingestion Architecture", [
    "📥 Multi-source Data Collection:",
    "   • Security logs (Windows, Linux, Network devices)",
    "   • API integrations (3rd-party security tools)",
    "   • Syslog and CEF formatted data",
    "   • Custom data connectors",
    "",
    "🔐 Data Processing:",
    "   • Real-time parsing and normalization",
    "   • Data enrichment and correlation",
    "   • Automatic threat classification"
])

# Slide 4: Data Flow
add_content_slide(prs, "Data Flow Pipeline", [
    "1️⃣  Data Collection",
    "   Sources → Collection Points → Message Queue",
    "",
    "2️⃣  Processing",
    "   Parsing → Enrichment → Correlation → Storage",
    "",
    "3️⃣  Indexing",
    "   Elasticsearch/Time-series database",
    "",
    "4️⃣  Visualization & Analytics",
    "   Dashboard → Search & Investigation → Alerts"
])

# Slide 5: Visualization Components
add_content_slide(prs, "Visualization Features", [
    "📈 Real-time Dashboards:",
    "   • Alert overview and severity distribution",
    "   • Security metrics and KPIs",
    "   • Threat intelligence summaries",
    "",
    "🔍 Interactive Search:",
    "   • SPL (Search Processing Language) queries",
    "   • Advanced filtering and time range selection",
    "   • Custom visualization types"
])

# Slide 6: Core Modules
add_content_slide(prs, "Core Modules", [
    "🔎 Search & Investigate",
    "   Query building, result visualization, timeline analysis",
    "",
    "⚠️ Alerts",
    "   Alert creation, management, and incident correlation",
    "",
    "🚨 Incidents",
    "   Incident tracking, investigation workflow, response actions",
    "",
    "📊 Analytics",
    "   Statistical analysis, trend detection, predictive insights"
])

# Slide 7: Search & Investigate Module
add_content_slide(prs, "Search & Investigate Module", [
    "🔍 Advanced Query Engine:",
    "   • SPL syntax support for complex searches",
    "   • Real-time query execution",
    "   • Result caching for performance",
    "",
    "📊 Result Visualization:",
    "   • Table views with sortable columns",
    "   • Timeline visualization",
    "   • Statistical breakdowns"
])

# Slide 8: Alerts Module
add_content_slide(prs, "Alerts Module", [
    "⚠️ Alert Management:",
    "   • Create custom alert rules",
    "   • Set severity levels (Critical, Warning, Info)",
    "   • Auto-correlation with related events",
    "",
    "🎯 Alert Features:",
    "   • Real-time alert triggering",
    "   • Alert deduplication",
    "   • Integration with incidents"
])

# Slide 9: Incidents Module
add_content_slide(prs, "Incidents Module", [
    "🚨 Incident Lifecycle:",
    "   • Creation from alerts or manual input",
    "   • Status tracking (Open, Investigating, Escalated, Resolved)",
    "   • Evidence collection and documentation",
    "",
    "👥 Collaboration:",
    "   • Team-based incident investigation",
    "   • Audit trail of all actions",
    "   • Automated response triggers"
])

# Slide 10: Key Features
add_content_slide(prs, "Key Features", [
    "🤖 AI-Native Analysis:",
    "   • Anomaly detection using ML models",
    "   • Behavioral analytics",
    "   • Threat intelligence matching",
    "",
    "⚡ Performance:",
    "   • Sub-second query response times",
    "   • Scalable architecture",
    "   • Real-time processing"
])

# Slide 11: Time Range Picker
add_content_slide(prs, "Time Range Selection", [
    "📅 Flexible Time Selection:",
    "   • Preset ranges (Last 24h, 7d, 30d)",
    "   • Custom date range picker",
    "   • Relative time expressions",
    "",
    "🔄 Data Retrieval:",
    "   • Optimized queries for time ranges",
    "   • Progressive data loading",
    "   • Historical data analysis"
])

# Slide 12: User Interface
add_content_slide(prs, "User Interface Design", [
    "🎨 Modern, Responsive Layout:",
    "   • Left panel: Navigation and filters",
    "   • Center panel: Main workspace",
    "   • Right panel: Contextual information",
    "",
    "⌨️ Keyboard Navigation:",
    "   • Quick search shortcuts",
    "   • Module switching",
    "   • Result filtering"
])

# Slide 13: Integration Points
add_content_slide(prs, "Integration & APIs", [
    "🔗 External Integrations:",
    "   • SIEM data sources (Splunk, ELK Stack)",
    "   • Threat intelligence feeds",
    "   • Security orchestration platforms",
    "",
    "📡 REST API:",
    "   • Query execution",
    "   • Alert management",
    "   • Incident creation and updates"
])

# Slide 14: Security & Compliance
add_content_slide(prs, "Security & Compliance", [
    "🔐 Data Security:",
    "   • Role-based access control (RBAC)",
    "   • Audit logging of all activities",
    "   • Encryption in transit and at rest",
    "",
    "📋 Compliance:",
    "   • GDPR, HIPAA, SOC2 compliance ready",
    "   • Regulatory reporting features"
])

# Slide 15: Future Enhancements
add_content_slide(prs, "Future Enhancements", [
    "🚀 Planned Features:",
    "   • Machine learning-based threat prediction",
    "   • Automated incident response playbooks",
    "   • Advanced correlation engine",
    "",
    "📊 Upcoming Modules:",
    "   • Vulnerability management",
    "   • Compliance monitoring",
    "   • Threat hunting automation"
])

# Slide 16: Conclusion
add_title_slide(prs, "Ready to Deploy", "Secure your enterprise with FACTOSIEM")

# Save presentation
prs.save('FACTOSIEM_Architecture_Guide.pptx')
print("✓ Presentation created: FACTOSIEM_Architecture_Guide.pptx")
